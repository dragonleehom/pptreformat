import streamlit as st 
import pandas as pd
from io import StringIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

from pptx.util import Cm,Pt
from pptx.dml.color import RGBColor
from PIL import Image
import math
import re
import os
import copy

def process_txt_format(unformated_text):
    print("-------------processing format of text",unformated_text)
    txt_remove_space=unformated_text.replace(" ","").replace("\n","").replace("\r","")
    print(txt_remove_space)
    splittxt = re.split('[(|)|（|）]',txt_remove_space)
    print(splittxt)
    for eachword in splittxt:
        print(eachword)
        if len(eachword)==0:
            splittxt.remove(eachword)
            continue
    return splittxt

def process_slide(slide_src):
    # 获取页面中的所有形状
    shapes = slide_src.shapes

    # 遍历所有形状
    for shape in shapes:
        # 判断形状类型是否为图片
        print(shape.name)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 处理图片
            #process_image(shape, slide_dst)
            if shape.width >Cm(0.1) and shape.height > Cm(0.1) and shape.height < prs_dst.slide_width:
                slide_dst = prs_dst.slides.add_slide(blank_slide_layout)

                imdata = shape.image.blob
                imagetype = shape.image.content_type
                typekey = imagetype.find('/') + 1
                imtype = imagetype[typekey:]
                image_file ="tmp."+imtype
                file_str=open(image_file,'wb')
                file_str.write(imdata)
                file_str.close()

                new_image_width=shape.width*3
                new_image_height=shape.height*3
                new_image_top=(prs_dst.slide_height/2)-(new_image_height/2)
                new_image_left=(prs_dst.slide_width/4)-(new_image_width/2)
                print("****slide size:",prs_dst.slide_width/360000,prs_dst.slide_height/360000)
                print("****add picture:",new_image_height/360000,new_image_width/360000,new_image_top/360000,new_image_left/360000)
                
                new_shape = slide_dst.shapes.add_picture(image_file,new_image_left,new_image_top,new_image_width,new_image_height)
 
                gap_left = 0
                gap_top  = 0
                latest_left=0
                latest_top=0
                gap_distance =0

                txBox = slide_dst.shapes.add_textbox(prs_dst.slide_width/2,0,prs_dst.slide_width/2,prs_dst.slide_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                tf.clear()
                
                #tf.fit_text(font_family='Calibri', max_size=72, bold=True, italic=False, font_file=None)
                tmp_target_text=""
                print("----start to match word ")
                for textShape in shapes:
                    print("--------",textShape.name,"----",textShape.shape_type,textShape)
                    if textShape.has_text_frame or textShape.shape_type==MSO_SHAPE_TYPE.TABLE:
                        
                        #if (textShape.left-shape.left)<0:
                        if textShape.has_text_frame:
                            print("------------text:",textShape.text_frame.text,textShape)
                            if textShape.text_frame.text.find("磨出好耳朵")!=-1:
                                print("------------continue:",textShape.left-shape.left,textShape.top-shape.top)
                                continue
                        print("----------------start to copy content----")
                        if len(tmp_target_text) == 0:
                            if textShape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                cellindex=0
                                for col in textShape.columns:
                                    eachcell = textShape.cell(0,col)
                                    if len(eachcell.text)==0:
                                        cellindex+=1
                                        continue
                                    latest_top = textShape.top
                                    latest_left = textShape.left+cellindex*textShape.column[0].width
                                    if gap_distance == 0:
                                        gap_distance = math.sqrt((latest_left-shape.left)**2+(latest_top-shape.top)**2)
                                        tmp_target_text = eachcard.text
                                    elif gap_distance >math.sqrt((latest_left-shape.left)**2+(latest_top-shape.top)**2):
                                        gap_distance = math.sqrt((latest_left-shape.left)**2+(latest_top-shape.top)**2)
                                        tmp_target_text = eachcard.text
                                    cellindex+=1
                            else:            
                                gap_left=textShape.left-shape.left
                                gap_top=abs(textShape.top - shape.top)
                                gap_distance = math.sqrt((textShape.left-shape.left)**2+(textShape.top-shape.top)**2)
                                tmp_target_text = textShape.text_frame.text
                            print("--------------------this is the first time----",gap_left,gap_top,gap_distance)
                        else:
                            if textShape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                cellindex=0
                                for col in textShape.columns:
                                    eachcell = textShape.cell(0,col)
                                    if len(eachcell.text)==0:
                                        cellindex+=1
                                        continue
                                    latest_top = textShape.top
                                    latest_left = textShape.left+cellindex*textShape.column[0].width
                                    if gap_distance >math.sqrt((latest_left-shape.left)**2+(latest_top-shape.top)**2):
                                        gap_distance = math.sqrt((latest_left-shape.left)**2+(latest_top-shape.top)**2)
                                        tmp_target_text = eachcard.text
                                    cellindex+=1
                            else:   
                                tmp_gap_left=textShape.left - shape.left
                                tmp_gap_top=abs(textShape.top - shape.top)
                                tmp_gap_distance=math.sqrt((textShape.left-shape.left)**2+(textShape.top-shape.top)**2)
                                print("--------------------tmp gap----",gap_left,tmp_gap_left,gap_top,tmp_gap_top,gap_distance,tmp_gap_distance)
                                if gap_distance>tmp_gap_distance:
                                    print("--------------------find closer one----",textShape.text_frame.text)
                                    if len(textShape.text_frame.text)>0:
                                        gap_left=tmp_gap_left
                                        gap_top=tmp_gap_top
                                        gap_distance=tmp_gap_distance
                                        tmp_target_text=textShape.text_frame.text
                
                #规范化字符
                cardwordlist= process_txt_format(tmp_target_text)

                i=0
                for eachcard in cardwordlist:
                    if i==0:
                        p = tf.add_paragraph()
                        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        p.text = eachcard
                        p.font.name="Calibri"
                        p.font.bold=True
                        p.font.size=Pt(80)
                    elif i==1:
                        p = tf.add_paragraph()
                        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        p.text = "(" + eachcard + ")"
                        p.font.name="微软雅黑"
                        p.font.bold=False
                        p.font.size=Pt(36)
                    elif i==2:
                        p = tf.add_paragraph()
                        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        p.text = eachcard
                        p.font.name="微软雅黑"
                        p.font.bold=False
                        p.font.size=Pt(36)
                    else:
                        print("---------Abnormal text-----",eachcard,i)
                    i+=1

                print("----end matching word")



         
        # elif shape.shape_type == MSO_SHAPE_TYPE.LINE
        #     shapes.element.remove(shape.element)



def app_head():
    #列举当前功能信息
    st.markdown("铭铭的英语学习卡片转换助手")
    st.write("已经实现功能")
    asis_data = {
        "时间": 
            ["2024-03-26",
        ],
        "功能": 
            ["读取固定ppt源文件生成固定目标ppt"
            "将源文件中的图片分别生成对应的A5尺寸的ppt页面"
            "将相关的文本拷贝到图片左侧，形成卡片",
        ]
    }
    asis_df = pd.DataFrame(asis_data)
    st.write(asis_df)

    #列举待实现的功能信息
    st.write("规划中的功能")
    rdmap_data = {
        "时间": 
            ["2024-03-31",
            "2024-04-20",
            
        ],
        "功能": 
            ["文件上传下载"
            "文字匹配更加智能，同时支持格式自动生成",
            "加入对pdf的支持",
        ]
    }
    rdmap_df = pd.DataFrame(rdmap_data)
    st.write(rdmap_df)


app_head()

uploaded_files=st.file_uploader("请上传需要调整格式的ppt文件，仅支持ppt/pptx,可同时上传多个文件",accept_multiple_files=True)
#uploaded_files.type=['ppt','pptx']
uploaded_file_name=""
for uploaded_file in uploaded_files:
    bytes_data = uploaded_file.read()
    uploaded_file_name=uploaded_file.name
    with open(uploaded_file_name, "wb") as f:
        f.write(uploaded_file.getbuffer())
    st.write("filename:", uploaded_file_name)
    #st.write(bytes_data)
    # 打开原始和目标 PPT
    prs_src = Presentation(uploaded_file_name)

    target_file_name = "转换后-" + uploaded_file_name
    # 新建目标 PPT
    prs_dst = Presentation()
    prs_dst.slide_height=Cm(14.8)
    prs_dst.slide_width=Cm(21.0)
    blank_slide_layout=prs_dst.slide_layouts[6]
    # 遍历原始 PPT 中的每一页
    for slide_src in prs_src.slides:
        # 处理当前页面的图片
        process_slide(slide_src)

    # 保存目标 PPT
    try:
        prs_dst.save(target_file_name)
    except FileExistsError:
        # 如果目标文件已存在，则直接覆盖
        prs_dst.save(target_file_name)

    with open(target_file_name, 'rb') as ff:
        target_file = ff
        st.download_button('下载转换后的pptx', target_file.read(),file_name=target_file_name,mime="pptx") 

#binary_contents = b'target.pptx'
#with open('myfile.zip', 'rb') as f:
#   st.download_button('Download target', f, file_name='target.ppx')


st.write("Here we are at the end of getting started with streamlit! Happy Streamlit-ing! :balloon:")
